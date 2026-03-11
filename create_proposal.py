from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

TOYOTA_RED = RGBColor(235, 10, 30)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(88, 89, 91)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = TOYOTA_RED
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, is_bullet=True):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TOYOTA_RED
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item if is_bullet else item
        p.font.size = Pt(22)
        p.font.color.rgb = BLACK
        p.space_after = Pt(16)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TOYOTA_RED
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5))
    tf = left_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TOYOTA_RED
    
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.add_paragraph()
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_after = Pt(12)
    
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5))
    tf = right_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TOYOTA_RED
    
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.add_paragraph()
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_after = Pt(12)
    
    return slide

def add_stats_slide(prs, title, stats_data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TOYOTA_RED
    
    col_width = Inches(3.8)
    start_x = Inches(0.8)
    gap = Inches(0.3)
    
    for i, (stat, desc) in enumerate(stats_data):
        x = start_x + i * (col_width + gap)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2), col_width, Inches(4))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = GRAY
        card.line.width = Pt(2)
        
        stat_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.5), col_width - Inches(0.4), Inches(1.5))
        tf = stat_box.text_frame
        p = tf.paragraphs[0]
        p.text = stat
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = TOYOTA_RED
        p.alignment = PP_ALIGN.CENTER
        
        desc_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(4.2), col_width - Inches(0.4), Inches(1.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_cta_slide(prs, title, items, cta_text):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLACK
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    content_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.333), Inches(3))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(14)
    
    cta_box = slide.shapes.add_textbox(Inches(4), Inches(5.8), Inches(5.333), Inches(1))
    tf = cta_box.text_frame
    p = tf.paragraphs[0]
    p.text = cta_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# Slide 1: Title
add_title_slide(prs, 
    "Toyota Mobility for Emergency",
    "Unified Emergency Response Platform for the Philippines")

# Slide 2: The Problem
add_content_slide(prs, "The Problem", [
    "Filipinos must memorize multiple emergency numbers that vary by city/municipality",
    "No unified system to track emergency status or response progress",
    "Delayed response times due to wrong number routing",
    "No visibility on responder location and ETA",
    "Current system: fragmented, voice-only, no tracking"
])

# Slide 3: Our Solution
add_two_column_slide(prs, "Our Solution",
    "What We Built", [
        "Single mobile app for Police, Ambulance, Fire emergencies",
        "Decentralized response - responders see emergencies directly in real-time",
        "Real-time location sharing and live ETA tracking",
        "First Aider Network - automatic alerts to nearby medical professionals",
        "Incident management - auto-merge duplicates, multi-unit support"
    ],
    "Target Users", [
        "General Filipino public",
        "Police (PNP), Ambulance, Fire (BFP)",
        "First aiders (nurses, doctors, trained volunteers)",
        "Dispatchers (monitoring/backup role)"
    ])

# Slide 4: Key Features
add_content_slide(prs, "Key Features", [
    "ONE-TAP EMERGENCY REPORTING - Report with auto-location, photos, and description in seconds",
    "DECENTRALIZED DISPATCH - Responders see emergencies directly (no dispatcher bottleneck)",
    "LIVE TRACKING - Reporter sees responder location and exact ETA in real-time",
    "FIRST AIDER ALERTS - System automatically notifies nearby medical professionals",
    "VERIFIED RESPONDERS - PhilID + supervisor approval for all responders"
])

# Slide 5: Impact & Metrics
add_stats_slide(prs, "Impact & Success Metrics", [
    ("20%", "Response Time Reduction"),
    ("500K", "Downloads in 6 Months"),
    ("4.5+", "App Store Rating"),
    ("<5%", "False Report Rate")
])

# Slide 6: Ask / Next Steps
add_cta_slide(prs, "Request Approval", [
    "Proceed with MVP development (9-12 month timeline)",
    "Secure initial budget: PHP 15-25M for MVP",
    "Initiate government partnership discussions with DILG",
    "Begin technical architecture planning",
    "Start pilot LGU selection process"
], "APPROVE TO PROCEED")

prs.save('Toyota_Mobility_Emergency_Proposal.pptx')
print("Presentation created: Toyota_Mobility_Emergency_Proposal.pptx")
